"""Append slides to mcp_slides.pptx: Data Security & Jira/Confluence Scenario Demo."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Colors ────────────────────────────────────────────────────────────────────
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_BG       = RGBColor(0xF8, 0xFA, 0xFC)
TITLE_COLOR    = RGBColor(0x1E, 0x29, 0x3B)
SUBTITLE_COLOR = RGBColor(0x33, 0x41, 0x55)
BODY_COLOR     = RGBColor(0x47, 0x55, 0x69)
MUTED_COLOR    = RGBColor(0x64, 0x74, 0x8B)
BOLD_COLOR     = RGBColor(0x1E, 0x29, 0x3B)
BLUE_ACCENT    = RGBColor(0x25, 0x63, 0xEB)
BLUE_DARK      = RGBColor(0x1E, 0x40, 0xAF)
BLUE_BG        = RGBColor(0xEF, 0xF6, 0xFF)
BLUE_CHIP_BG   = RGBColor(0xDB, 0xEA, 0xFE)
BLUE_CHIP_BD   = RGBColor(0x93, 0xC5, 0xFD)
GREEN_ACCENT   = RGBColor(0x16, 0xA3, 0x4A)
GREEN_DARK     = RGBColor(0x16, 0x65, 0x34)
GREEN_BG       = RGBColor(0xF0, 0xFD, 0xF4)
GREEN_CHIP_BG  = RGBColor(0xDC, 0xFC, 0xE7)
GREEN_CHIP_BD  = RGBColor(0x86, 0xEF, 0xAC)
PURPLE_ACCENT  = RGBColor(0xA8, 0x55, 0xF7)
PURPLE_DARK    = RGBColor(0x6B, 0x21, 0xA8)
PURPLE_BG      = RGBColor(0xFA, 0xF5, 0xFF)
GRAY_BD        = RGBColor(0x94, 0xA3, 0xB8)
GRAY_LIGHT     = RGBColor(0xCB, 0xD5, 0xE1)
YELLOW_BG      = RGBColor(0xFE, 0xF9, 0xC3)
YELLOW_BD      = RGBColor(0xFA, 0xCC, 0x15)
YELLOW_TEXT     = RGBColor(0x85, 0x4D, 0x0E)
ORANGE_ACCENT  = RGBColor(0xEA, 0x58, 0x0C)
ORANGE_DARK    = RGBColor(0x9A, 0x34, 0x12)
ORANGE_BG      = RGBColor(0xFF, 0xF7, 0xED)
ORANGE_CHIP_BG = RGBColor(0xFE, 0xD7, 0xAA)
ORANGE_CHIP_BD = RGBColor(0xFD, 0xBA, 0x74)
RED_ACCENT     = RGBColor(0xDC, 0x26, 0x26)
RED_DARK       = RGBColor(0x99, 0x1B, 0x1B)
RED_BG         = RGBColor(0xFE, 0xF2, 0xF2)
RED_CHIP_BG    = RGBColor(0xFE, 0xCD, 0xCD)
RED_CHIP_BD    = RGBColor(0xFC, 0xA5, 0xA5)
CODE_BG        = RGBColor(0xF1, 0xF5, 0xF9)
CODE_BD        = RGBColor(0xCB, 0xD5, 0xE1)
INDIGO_BG      = RGBColor(0xE8, 0xEA, 0xF6)
INDIGO_BD      = RGBColor(0x3F, 0x51, 0xB5)
INDIGO_DARK    = RGBColor(0x28, 0x35, 0x93)

FONT_NAME = "Segoe UI"

# Open existing deck
pptx_path = os.path.join(os.path.dirname(__file__), "mcp_slides.pptx")
prs = Presentation(pptx_path)

# ── Helpers ──────────────────────────────────────────────────────────────────

def make_slide_bg(s):
    bg = s.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = SLIDE_BG


def add_textbox(s, left, top, width, height, text, font_size=11, bold=False,
                color=BODY_COLOR, alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def add_rounded_rect(s, left, top, width, height, fill_color, border_color,
                     border_width=Pt(1.5)):
    shape = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    shape.adjustments[0] = 0.08
    return shape


def add_rich_textbox(s, left, top, width, height):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def add_connector(s, x1, y1, x2, y2, color, width=Pt(1.5), dashed=False):
    cx = s.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    cx.line.color.rgb = color
    cx.line.width = width
    if dashed:
        cx.line.dash_style = 4
    return cx


def add_arrow_line(s, x1, y1, x2, y2, color, width=Pt(1.5), dashed=False):
    cx = add_connector(s, x1, y1, x2, y2, color, width, dashed)
    ln = cx.line._ln
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return cx


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Data Security & Jira/Confluence Scenario Demo
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
make_slide_bg(slide)

# ── Title ──
add_textbox(slide, 0.5, 0.2, 12, 0.5,
            "Data Security & Jira / Confluence Scenario Demo",
            font_size=24, bold=True, color=TITLE_COLOR)
add_textbox(slide, 0.5, 0.7, 12, 0.35,
            "OBO authentication + tool-level authorization — per-user access control via MCP middleware",
            font_size=13, color=MUTED_COLOR)

# ═════════════════════════════════════════════════════════════════════
# LEFT SIDE — Security Architecture (text / cards)
# ═════════════════════════════════════════════════════════════════════

LX = 0.3

# -- Card: End-to-End Flow --
card_flow = add_rounded_rect(slide, LX, 1.2, 6.0, 2.6, WHITE, GRAY_LIGHT, Pt(1))
card_flow.adjustments[0] = 0.04

add_textbox(slide, LX + 0.2, 1.3, 5.5, 0.35,
            "End-to-End OBO Security Flow", font_size=16, bold=True, color=BLUE_ACCENT)

flow_steps = [
    ("1. MSAL Login", " — User authenticates via React SPA → Entra ID token"),
    ("2. OBO Exchange", " — FastAPI backend exchanges token for Foundry + Atlassian scopes"),
    ("3. Agent Invocation", " — Foundry Agent receives request with OBO token attached"),
    ("4. MCP Authentication", " — Custom MCP server validates Atlassian token via /accessible-resources"),
    ("5. Tool Authorization", " — AccessListMiddleware checks user sub → filters visible tools"),
    ("6. Scoped Execution", " — Only permitted tools execute; denied calls raise AuthorizationError"),
]

y = 1.7
for label, desc in flow_steps:
    tf = add_rich_textbox(slide, LX + 0.25, y, 5.4, 0.33)
    p = tf.paragraphs[0]
    rb = p.add_run(); rb.text = label
    rb.font.size = Pt(11); rb.font.bold = True
    rb.font.color.rgb = BOLD_COLOR; rb.font.name = FONT_NAME
    rd = p.add_run(); rd.text = desc
    rd.font.size = Pt(10); rd.font.color.rgb = BODY_COLOR; rd.font.name = FONT_NAME
    y += 0.29

# -- Card: Key Security Concepts --
card_sec = add_rounded_rect(slide, LX, 4.0, 6.0, 3.2, WHITE, GRAY_LIGHT, Pt(1))
card_sec.adjustments[0] = 0.04

add_textbox(slide, LX + 0.2, 4.1, 5.5, 0.35,
            "Key Security Concepts", font_size=16, bold=True, color=GREEN_ACCENT)

concepts = [
    ("On-Behalf-Of (OBO)", "User identity preserved end-to-end; no service-to-service tokens bypass user context"),
    ("Tool-Level Authorization", "AccessListMiddleware filters tools per user sub. Both list and call are guarded"),
    ("user_access_list.json", "Declarative per-user config: { sub → [tools] }. Hot-reloaded on change, no restart"),
    ("Premium vs Freemium", "Premium: Jira + Confluence + Fabric. Freemium: Jira + Fabric only"),
    ("Fail Closed", "No token / missing sub / tool not in list → AuthorizationError. Hidden tools can't be called directly"),
]

y = 4.5
for label, desc in concepts:
    tf = add_rich_textbox(slide, LX + 0.25, y, 5.4, 0.52)
    p = tf.paragraphs[0]
    p.space_after = Pt(2)
    rb = p.add_run(); rb.text = "• " + label + ": "
    rb.font.size = Pt(11); rb.font.bold = True
    rb.font.color.rgb = BOLD_COLOR; rb.font.name = FONT_NAME
    rd = p.add_run(); rd.text = desc
    rd.font.size = Pt(10); rd.font.color.rgb = BODY_COLOR; rd.font.name = FONT_NAME
    y += 0.52

# ═════════════════════════════════════════════════════════════════════
# RIGHT SIDE — Architecture Diagram (native shapes)
# ═════════════════════════════════════════════════════════════════════

RX = 6.7

diag_card = add_rounded_rect(slide, RX - 0.1, 1.2, 6.5, 6.0, WHITE, GRAY_LIGHT, Pt(1))
diag_card.adjustments[0] = 0.04

add_textbox(slide, RX, 1.3, 6.2, 0.35,
            "Architecture — OBO Token Flow",
            font_size=14, bold=True, color=TITLE_COLOR, alignment=PP_ALIGN.CENTER)

# ── Row 1: Users → SPA → Backend → Foundry ──

# User A (Premium)
ua_x, ua_y = RX + 0.1, 1.85
add_rounded_rect(slide, ua_x, ua_y, 1.05, 0.7, GREEN_BG, GREEN_ACCENT, Pt(1.5))
add_textbox(slide, ua_x, ua_y + 0.06, 1.05, 0.22, "User A",
            font_size=9, bold=True, color=GREEN_DARK, alignment=PP_ALIGN.CENTER)
add_textbox(slide, ua_x, ua_y + 0.28, 1.05, 0.2, "Premium",
            font_size=8, color=GREEN_DARK, alignment=PP_ALIGN.CENTER)
add_textbox(slide, ua_x, ua_y + 0.46, 1.05, 0.2, "Jira+Confluence+Fabric",
            font_size=7, color=BODY_COLOR, alignment=PP_ALIGN.CENTER)

# User B (Freemium)
ub_x, ub_y = RX + 0.1, 2.7
add_rounded_rect(slide, ub_x, ub_y, 1.05, 0.7, ORANGE_BG, ORANGE_ACCENT, Pt(1.5))
add_textbox(slide, ub_x, ub_y + 0.06, 1.05, 0.22, "User B",
            font_size=9, bold=True, color=ORANGE_DARK, alignment=PP_ALIGN.CENTER)
add_textbox(slide, ub_x, ub_y + 0.28, 1.05, 0.2, "Freemium",
            font_size=8, color=ORANGE_DARK, alignment=PP_ALIGN.CENTER)
add_textbox(slide, ub_x, ub_y + 0.46, 1.05, 0.2, "Jira+Fabric only",
            font_size=7, color=BODY_COLOR, alignment=PP_ALIGN.CENTER)

# React SPA
spa_x, spa_y = RX + 1.5, 1.95
spa_w, spa_h = 1.15, 1.35
add_rounded_rect(slide, spa_x, spa_y, spa_w, spa_h, BLUE_BG, BLUE_ACCENT, Pt(1.5))
add_textbox(slide, spa_x, spa_y + 0.05, spa_w, 0.25, "React SPA",
            font_size=10, bold=True, color=BLUE_DARK, alignment=PP_ALIGN.CENTER)
for i, txt in enumerate(["MSAL Auth", "Entra ID", "OBO Request"]):
    chip = add_rounded_rect(slide, spa_x + 0.08, spa_y + 0.35 + i * 0.3, 0.99, 0.25, BLUE_CHIP_BG, BLUE_CHIP_BD, Pt(1))
    chip.adjustments[0] = 0.15
    add_textbox(slide, spa_x + 0.08, spa_y + 0.37 + i * 0.3, 0.99, 0.23, txt,
                font_size=8, color=BLUE_DARK, alignment=PP_ALIGN.CENTER)

# Arrows Users → SPA
add_arrow_line(slide, ua_x + 1.05, ua_y + 0.35, spa_x, spa_y + 0.4, GRAY_BD, Pt(1))
add_arrow_line(slide, ub_x + 1.05, ub_y + 0.35, spa_x, spa_y + 0.95, GRAY_BD, Pt(1))

# FastAPI Backend
be_x, be_y = RX + 3.0, 1.95
be_w, be_h = 1.15, 1.35
add_rounded_rect(slide, be_x, be_y, be_w, be_h, PURPLE_BG, PURPLE_ACCENT, Pt(1.5))
add_textbox(slide, be_x, be_y + 0.05, be_w, 0.25, "FastAPI",
            font_size=10, bold=True, color=PURPLE_DARK, alignment=PP_ALIGN.CENTER)
for i, txt in enumerate(["JWT Validate", "OBO Exchange", "Scope Mgmt"]):
    chip_c = RGBColor(0xF3, 0xE8, 0xFF)
    chip_b = RGBColor(0xD8, 0xB4, 0xFE)
    chip = add_rounded_rect(slide, be_x + 0.08, be_y + 0.35 + i * 0.3, 0.99, 0.25, chip_c, chip_b, Pt(1))
    chip.adjustments[0] = 0.15
    add_textbox(slide, be_x + 0.08, be_y + 0.37 + i * 0.3, 0.99, 0.23, txt,
                font_size=8, color=PURPLE_DARK, alignment=PP_ALIGN.CENTER)

# Arrow SPA → Backend
add_arrow_line(slide, spa_x + spa_w, spa_y + spa_h / 2, be_x, be_y + be_h / 2, BLUE_ACCENT, Pt(1.5))
add_textbox(slide, spa_x + spa_w + 0.02, spa_y + spa_h / 2 - 0.2, 0.3, 0.18, "Bearer",
            font_size=7, bold=True, color=BLUE_ACCENT, alignment=PP_ALIGN.CENTER)

# Foundry Agent
fa_x, fa_y = RX + 4.5, 1.85
fa_w, fa_h = 1.5, 1.55
add_rounded_rect(slide, fa_x, fa_y, fa_w, fa_h, RED_BG, RED_ACCENT, Pt(1.5))
add_textbox(slide, fa_x, fa_y + 0.05, fa_w, 0.25, "AI Foundry Agent",
            font_size=10, bold=True, color=RED_DARK, alignment=PP_ALIGN.CENTER)
for i, txt in enumerate(["gpt-4.1-mini", "MCP Tool Conn", "Fabric Tool Conn"]):
    chip = add_rounded_rect(slide, fa_x + 0.08, fa_y + 0.35 + i * 0.3, 1.34, 0.25, RED_CHIP_BG, RED_CHIP_BD, Pt(1))
    chip.adjustments[0] = 0.15
    add_textbox(slide, fa_x + 0.08, fa_y + 0.37 + i * 0.3, 1.34, 0.23, txt,
                font_size=8, color=RED_DARK, alignment=PP_ALIGN.CENTER)

# OBO token callout
obo_y = fa_y + 1.15
add_rounded_rect(slide, fa_x + 0.05, obo_y, 1.4, 0.3, YELLOW_BG, YELLOW_BD, Pt(1))
add_textbox(slide, fa_x + 0.05, obo_y + 0.03, 1.4, 0.25, "OBO token → MCP",
            font_size=8, bold=True, color=YELLOW_TEXT, alignment=PP_ALIGN.CENTER)

# Arrow Backend → Foundry
add_arrow_line(slide, be_x + be_w, be_y + be_h / 2, fa_x, fa_y + fa_h / 2, PURPLE_ACCENT, Pt(1.5))
add_textbox(slide, be_x + be_w + 0.02, be_y + be_h / 2 - 0.2, 0.3, 0.18, "OBO",
            font_size=7, bold=True, color=PURPLE_DARK, alignment=PP_ALIGN.CENTER)

# ── Row 2: MCP Server Security Layers ──

sec_y = 3.65

# Section label bar
add_rounded_rect(slide, RX + 0.05, sec_y, 6.2, 0.35, INDIGO_BG, INDIGO_BD, Pt(1))
add_textbox(slide, RX + 0.1, sec_y + 0.04, 6.1, 0.28,
            "Custom Atlassian MCP Server — Security Layers",
            font_size=10, bold=True, color=INDIGO_DARK, alignment=PP_ALIGN.CENTER)

ly = sec_y + 0.5

# Layer 1: Authentication
l1_x = RX + 0.1
l1_w, l1_h = 1.95, 1.55
add_rounded_rect(slide, l1_x, ly, l1_w, l1_h, ORANGE_BG, ORANGE_ACCENT, Pt(1.5))
add_textbox(slide, l1_x, ly + 0.05, l1_w, 0.22, "① Authentication",
            font_size=10, bold=True, color=ORANGE_DARK, alignment=PP_ALIGN.CENTER)

auth_items = ["DebugTokenVerifier", "Atlassian API check", "JWT → extract sub"]
for i, txt in enumerate(auth_items):
    chip = add_rounded_rect(slide, l1_x + 0.08, ly + 0.32 + i * 0.32, l1_w - 0.16, 0.27, ORANGE_CHIP_BG, ORANGE_CHIP_BD, Pt(1))
    chip.adjustments[0] = 0.15
    add_textbox(slide, l1_x + 0.08, ly + 0.34 + i * 0.32, l1_w - 0.16, 0.25, txt,
                font_size=8, color=ORANGE_DARK, alignment=PP_ALIGN.CENTER)

# Fail box
fail_y = ly + 1.25
add_rounded_rect(slide, l1_x + 0.15, fail_y, l1_w - 0.3, 0.22, RED_CHIP_BG, RED_ACCENT, Pt(1))
add_textbox(slide, l1_x + 0.15, fail_y + 0.02, l1_w - 0.3, 0.2, "✗ Invalid → 401",
            font_size=7, bold=True, color=RED_DARK, alignment=PP_ALIGN.CENTER)

# Arrow Layer 1 → Layer 2
add_arrow_line(slide, l1_x + l1_w, ly + l1_h / 2, l1_x + l1_w + 0.15, ly + l1_h / 2, GREEN_ACCENT, Pt(1.5))

# Layer 2: Authorization
l2_x = RX + 2.2
l2_w, l2_h = 2.0, 1.55
add_rounded_rect(slide, l2_x, ly, l2_w, l2_h, GREEN_BG, GREEN_ACCENT, Pt(1.5))
add_textbox(slide, l2_x, ly + 0.05, l2_w, 0.22, "② AccessListMiddleware",
            font_size=10, bold=True, color=GREEN_DARK, alignment=PP_ALIGN.CENTER)

auth_z_items = ["on_list_tools()", "on_call_tool()", "user_access_list.json"]
for i, txt in enumerate(auth_z_items):
    chip = add_rounded_rect(slide, l2_x + 0.08, ly + 0.32 + i * 0.32, l2_w - 0.16, 0.27, GREEN_CHIP_BG, GREEN_CHIP_BD, Pt(1))
    chip.adjustments[0] = 0.15
    add_textbox(slide, l2_x + 0.08, ly + 0.34 + i * 0.32, l2_w - 0.16, 0.25, txt,
                font_size=8, color=GREEN_DARK, alignment=PP_ALIGN.CENTER)

# Denied box
add_rounded_rect(slide, l2_x + 0.15, fail_y, l2_w - 0.3, 0.22, RED_CHIP_BG, RED_ACCENT, Pt(1))
add_textbox(slide, l2_x + 0.15, fail_y + 0.02, l2_w - 0.3, 0.2, "✗ AuthorizationError",
            font_size=7, bold=True, color=RED_DARK, alignment=PP_ALIGN.CENTER)

# Arrow Layer 2 → Layer 3
add_arrow_line(slide, l2_x + l2_w, ly + l2_h / 2, l2_x + l2_w + 0.15, ly + l2_h / 2, GREEN_ACCENT, Pt(1.5))

# Layer 3: Tool Execution
l3_x = RX + 4.35
l3_w, l3_h = 1.85, 1.55
add_rounded_rect(slide, l3_x, ly, l3_w, l3_h, PURPLE_BG, PURPLE_ACCENT, Pt(1.5))
add_textbox(slide, l3_x, ly + 0.05, l3_w, 0.22, "③ Tools",
            font_size=10, bold=True, color=PURPLE_DARK, alignment=PP_ALIGN.CENTER)

tools = ["jira_list_issues", "confluence_search", "Fabric Data Agent"]
for i, txt in enumerate(tools):
    chip_c = RGBColor(0xF3, 0xE8, 0xFF)
    chip_b = RGBColor(0xD8, 0xB4, 0xFE)
    chip = add_rounded_rect(slide, l3_x + 0.08, ly + 0.32 + i * 0.32, l3_w - 0.16, 0.27, chip_c, chip_b, Pt(1))
    chip.adjustments[0] = 0.15
    add_textbox(slide, l3_x + 0.08, ly + 0.34 + i * 0.32, l3_w - 0.16, 0.25, txt,
                font_size=8, color=PURPLE_DARK, alignment=PP_ALIGN.CENTER)

# Atlassian APIs box
api_y = ly + 1.25
add_rounded_rect(slide, l3_x + 0.1, api_y, l3_w - 0.2, 0.22, BLUE_CHIP_BG, BLUE_CHIP_BD, Pt(1))
add_textbox(slide, l3_x + 0.1, api_y + 0.02, l3_w - 0.2, 0.2, "→ Atlassian Cloud APIs",
            font_size=7, bold=True, color=BLUE_DARK, alignment=PP_ALIGN.CENTER)

# ── Row 3: Demo Scenario (Premium vs Freemium) ──

demo_y = 5.8

# Premium callout
add_rounded_rect(slide, RX + 0.1, demo_y, 3.0, 1.25, GREEN_BG, GREEN_ACCENT, Pt(1.5))
add_textbox(slide, RX + 0.2, demo_y + 0.05, 2.7, 0.25,
            "User A — Premium", font_size=11, bold=True, color=GREEN_DARK)

prem_tools = [
    ("✓ jira_list_issues", GREEN_CHIP_BG, GREEN_CHIP_BD, GREEN_DARK),
    ("✓ confluence_search", GREEN_CHIP_BG, GREEN_CHIP_BD, GREEN_DARK),
    ("✓ Fabric Data Agent", GREEN_CHIP_BG, GREEN_CHIP_BD, GREEN_DARK),
]
for i, (txt, bg, bd, tc) in enumerate(prem_tools):
    chip = add_rounded_rect(slide, RX + 0.2 + i * 0.95, demo_y + 0.35, 0.88, 0.25, bg, bd, Pt(1))
    chip.adjustments[0] = 0.15
    add_textbox(slide, RX + 0.2 + i * 0.95, demo_y + 0.37, 0.88, 0.23, txt,
                font_size=7, color=tc, alignment=PP_ALIGN.CENTER)

add_textbox(slide, RX + 0.2, demo_y + 0.68, 2.7, 0.2,
            '"Show Jira issues + search Confluence"',
            font_size=8, color=BODY_COLOR, font_name="Segoe UI")
add_textbox(slide, RX + 0.2, demo_y + 0.88, 2.7, 0.2,
            "→ All 3 tools available, full access",
            font_size=8, bold=True, color=GREEN_DARK)

# Freemium callout
add_rounded_rect(slide, RX + 3.25, demo_y, 3.0, 1.25, ORANGE_BG, ORANGE_ACCENT, Pt(1.5))
add_textbox(slide, RX + 3.35, demo_y + 0.05, 2.7, 0.25,
            "User B — Freemium", font_size=11, bold=True, color=ORANGE_DARK)

freem_tools = [
    ("✓ jira_list_issues", GREEN_CHIP_BG, GREEN_CHIP_BD, GREEN_DARK),
    ("✗ confluence_search", RED_CHIP_BG, RED_CHIP_BD, RED_DARK),
    ("✓ Fabric Data Agent", GREEN_CHIP_BG, GREEN_CHIP_BD, GREEN_DARK),
]
for i, (txt, bg, bd, tc) in enumerate(freem_tools):
    chip = add_rounded_rect(slide, RX + 3.35 + i * 0.95, demo_y + 0.35, 0.88, 0.25, bg, bd, Pt(1))
    chip.adjustments[0] = 0.15
    add_textbox(slide, RX + 3.35 + i * 0.95, demo_y + 0.37, 0.88, 0.23, txt,
                font_size=7, color=tc, alignment=PP_ALIGN.CENTER)

add_textbox(slide, RX + 3.35, demo_y + 0.68, 2.7, 0.2,
            '"Search Confluence for onboarding docs"',
            font_size=8, color=BODY_COLOR, font_name="Segoe UI")
add_textbox(slide, RX + 3.35, demo_y + 0.88, 2.7, 0.2,
            "→ AuthorizationError: Not authorized",
            font_size=8, bold=True, color=RED_DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Detailed Architecture Diagram (SVG embedded as image)
# ══════════════════════════════════════════════════════════════════════════════

s2 = prs.slides.add_slide(prs.slide_layouts[6])
make_slide_bg(s2)

add_textbox(s2, 0.5, 0.2, 12, 0.5,
            "Data Security Deep Dive — MCP Server Internals",
            font_size=24, bold=True, color=TITLE_COLOR)
add_textbox(s2, 0.5, 0.7, 12, 0.35,
            "Authentication → Authorization → Tool Execution — three security layers protect every request",
            font_size=13, color=MUTED_COLOR)

# ── LEFT: Three security layers detail ──

LX2 = 0.3

# Layer 1 Card
l1_card = add_rounded_rect(s2, LX2, 1.2, 4.0, 1.85, WHITE, ORANGE_ACCENT, Pt(1.5))
l1_card.adjustments[0] = 0.04
add_textbox(s2, LX2 + 0.2, 1.3, 3.5, 0.3,
            "① Authentication Layer", font_size=15, bold=True, color=ORANGE_DARK)

auth_details = [
    ("Token Validation", "DebugTokenVerifier calls Atlassian /accessible-resources"),
    ("JWT Decode", "Extract user sub claim (unverified) for identity mapping"),
    ("Cloud ID Cache", "Memoize Atlassian cloud IDs per user sub for performance"),
    ("Fail Fast", "Invalid token → 401 before reaching any tool logic"),
]
y = 1.65
for label, desc in auth_details:
    tf = add_rich_textbox(s2, LX2 + 0.25, y, 3.5, 0.3)
    p = tf.paragraphs[0]
    rb = p.add_run(); rb.text = label + ": "
    rb.font.size = Pt(10); rb.font.bold = True
    rb.font.color.rgb = ORANGE_DARK; rb.font.name = FONT_NAME
    rd = p.add_run(); rd.text = desc
    rd.font.size = Pt(9); rd.font.color.rgb = BODY_COLOR; rd.font.name = FONT_NAME
    y += 0.28

# Layer 2 Card
l2_card = add_rounded_rect(s2, LX2, 3.2, 4.0, 2.0, WHITE, GREEN_ACCENT, Pt(1.5))
l2_card.adjustments[0] = 0.04
add_textbox(s2, LX2 + 0.2, 3.3, 3.5, 0.3,
            "② AccessListMiddleware", font_size=15, bold=True, color=GREEN_DARK)

mw_details = [
    ("on_list_tools()", "Decode JWT → load user_access_list.json → return only allowed tools"),
    ("on_call_tool()", "Verify tool name in user's allowed set; raise AuthorizationError if denied"),
    ("Hot Reload", "Watches file mtime; auto-reloads permissions without server restart"),
    ("Wildcard", 'Grant all tools with "*" in allowed_tools list'),
    ("Thread Safe", "File loading protected by threading lock for concurrent requests"),
]
y = 3.65
for label, desc in mw_details:
    tf = add_rich_textbox(s2, LX2 + 0.25, y, 3.5, 0.3)
    p = tf.paragraphs[0]
    rb = p.add_run(); rb.text = label + ": "
    rb.font.size = Pt(10); rb.font.bold = True
    rb.font.color.rgb = GREEN_DARK; rb.font.name = FONT_NAME
    rd = p.add_run(); rd.text = desc
    rd.font.size = Pt(9); rd.font.color.rgb = BODY_COLOR; rd.font.name = FONT_NAME
    y += 0.28

# Layer 3 Card
l3_card = add_rounded_rect(s2, LX2, 5.35, 4.0, 1.85, WHITE, PURPLE_ACCENT, Pt(1.5))
l3_card.adjustments[0] = 0.04
add_textbox(s2, LX2 + 0.2, 5.45, 3.5, 0.3,
            "③ Tool Execution Layer", font_size=15, bold=True, color=PURPLE_DARK)

tool_details = [
    ("jira_list_issues", "Resolve Jira cloud_id → POST /search/jql with user's OBO token"),
    ("confluence_search_pages", "Resolve Confluence site → CQL search via /wiki/rest/api/search"),
    ("Cloud ID Resolution", "Probes accessible resources to find correct Atlassian site per user"),
    ("Fabric Data Agent", "Microsoft-managed MCP tool with Row-Level Security on Lakehouse data"),
]
y = 5.8
for label, desc in tool_details:
    tf = add_rich_textbox(s2, LX2 + 0.25, y, 3.5, 0.3)
    p = tf.paragraphs[0]
    rb = p.add_run(); rb.text = label + ": "
    rb.font.size = Pt(10); rb.font.bold = True
    rb.font.color.rgb = PURPLE_DARK; rb.font.name = FONT_NAME
    rd = p.add_run(); rd.text = desc
    rd.font.size = Pt(9); rd.font.color.rgb = BODY_COLOR; rd.font.name = FONT_NAME
    y += 0.28

# ── RIGHT: Architecture flow diagram ──

RX2 = 4.6
diag2 = add_rounded_rect(s2, RX2, 1.2, 8.4, 5.95, WHITE, GRAY_LIGHT, Pt(1))
diag2.adjustments[0] = 0.04

add_textbox(s2, RX2 + 0.1, 1.3, 8.1, 0.3,
            "Request Flow — MCP Client → Auth → Middleware → Tools → Atlassian APIs",
            font_size=12, bold=True, color=TITLE_COLOR, alignment=PP_ALIGN.CENTER)

# MCP Client
mc_x, mc_y = RX2 + 0.3, 1.85
mc_w, mc_h = 1.5, 0.8
add_rounded_rect(s2, mc_x, mc_y, mc_w, mc_h, BLUE_BG, BLUE_ACCENT, Pt(1.5))
add_textbox(s2, mc_x, mc_y + 0.08, mc_w, 0.25, "MCP Client",
            font_size=11, bold=True, color=BLUE_DARK, alignment=PP_ALIGN.CENTER)
add_textbox(s2, mc_x, mc_y + 0.35, mc_w, 0.2, "Foundry Agent / AI",
            font_size=9, color=MUTED_COLOR, alignment=PP_ALIGN.CENTER)
add_textbox(s2, mc_x, mc_y + 0.55, mc_w, 0.2, "Bearer Token attached",
            font_size=8, color=BLUE_DARK, alignment=PP_ALIGN.CENTER)

# Arrow → Auth
add_arrow_line(s2, mc_x + mc_w, mc_y + mc_h / 2, mc_x + mc_w + 0.3, mc_y + mc_h / 2, GRAY_BD, Pt(1.5))

# Auth box
auth_x = mc_x + mc_w + 0.35
auth_w, auth_h = 1.7, 0.8
add_rounded_rect(s2, auth_x, mc_y, auth_w, auth_h, ORANGE_BG, ORANGE_ACCENT, Pt(1.5))
add_textbox(s2, auth_x, mc_y + 0.08, auth_w, 0.25, "① Auth Layer",
            font_size=10, bold=True, color=ORANGE_DARK, alignment=PP_ALIGN.CENTER)
add_textbox(s2, auth_x, mc_y + 0.35, auth_w, 0.2, "Validate Atlassian token",
            font_size=8, color=BODY_COLOR, alignment=PP_ALIGN.CENTER)
add_textbox(s2, auth_x, mc_y + 0.55, auth_w, 0.2, "Decode JWT → sub claim",
            font_size=8, color=BODY_COLOR, alignment=PP_ALIGN.CENTER)

# Arrow → Middleware
add_arrow_line(s2, auth_x + auth_w, mc_y + mc_h / 2, auth_x + auth_w + 0.3, mc_y + mc_h / 2, GREEN_ACCENT, Pt(1.5))

# Middleware box
mw_x = auth_x + auth_w + 0.35
mw_w, mw_h = 1.7, 0.8
add_rounded_rect(s2, mw_x, mc_y, mw_w, mw_h, GREEN_BG, GREEN_ACCENT, Pt(1.5))
add_textbox(s2, mw_x, mc_y + 0.08, mw_w, 0.25, "② Middleware",
            font_size=10, bold=True, color=GREEN_DARK, alignment=PP_ALIGN.CENTER)
add_textbox(s2, mw_x, mc_y + 0.35, mw_w, 0.2, "Filter tools per user",
            font_size=8, color=BODY_COLOR, alignment=PP_ALIGN.CENTER)
add_textbox(s2, mw_x, mc_y + 0.55, mw_w, 0.2, "Guard tool calls",
            font_size=8, color=BODY_COLOR, alignment=PP_ALIGN.CENTER)

# Arrow → Tools
add_arrow_line(s2, mw_x + mw_w, mc_y + mc_h / 2, mw_x + mw_w + 0.3, mc_y + mc_h / 2, PURPLE_ACCENT, Pt(1.5))

# Tools box
tl_x = mw_x + mw_w + 0.35
tl_w, tl_h = 1.6, 0.8
add_rounded_rect(s2, tl_x, mc_y, tl_w, tl_h, PURPLE_BG, PURPLE_ACCENT, Pt(1.5))
add_textbox(s2, tl_x, mc_y + 0.08, tl_w, 0.25, "③ Tools",
            font_size=10, bold=True, color=PURPLE_DARK, alignment=PP_ALIGN.CENTER)
add_textbox(s2, tl_x, mc_y + 0.35, tl_w, 0.2, "jira / confluence",
            font_size=8, color=BODY_COLOR, alignment=PP_ALIGN.CENTER)
add_textbox(s2, tl_x, mc_y + 0.55, tl_w, 0.2, "→ Atlassian APIs",
            font_size=8, color=BODY_COLOR, alignment=PP_ALIGN.CENTER)

# Denied path (from middleware down)
deny_y = mc_y + mc_h + 0.2
add_arrow_line(s2, mw_x + mw_w / 2, mc_y + mc_h, mw_x + mw_w / 2, deny_y + 0.05, RED_ACCENT, Pt(1.5))
deny_box = add_rounded_rect(s2, mw_x - 0.1, deny_y + 0.1, mw_w + 0.2, 0.45, RED_BG, RED_ACCENT, Pt(1.5))
add_textbox(s2, mw_x - 0.05, deny_y + 0.13, mw_w + 0.1, 0.2, "✗ AuthorizationError",
            font_size=9, bold=True, color=RED_DARK, alignment=PP_ALIGN.CENTER)
add_textbox(s2, mw_x - 0.05, deny_y + 0.32, mw_w + 0.1, 0.18, "tool not in user's allowed set",
            font_size=8, color=BODY_COLOR, alignment=PP_ALIGN.CENTER)

# ── user_access_list.json visual ──
json_y = 3.2
json_card = add_rounded_rect(s2, RX2 + 0.2, json_y, 7.9, 1.6, CODE_BG, CODE_BD, Pt(1))
json_card.adjustments[0] = 0.03
add_textbox(s2, RX2 + 0.35, json_y + 0.05, 3.0, 0.25,
            "user_access_list.json", font_size=12, bold=True, color=TITLE_COLOR)
add_textbox(s2, RX2 + 3.5, json_y + 0.08, 4.0, 0.22,
            "(hot-reloaded on change — no server restart needed)",
            font_size=9, color=MUTED_COLOR)

# Premium user entry
add_rounded_rect(s2, RX2 + 0.35, json_y + 0.38, 3.7, 1.1, GREEN_BG, GREEN_ACCENT, Pt(1))
add_textbox(s2, RX2 + 0.5, json_y + 0.42, 3.3, 0.22,
            "User A (sub: 70121:9f08...)", font_size=10, bold=True, color=GREEN_DARK)
json_text_a = '[  "confluence_search_pages",  "jira_list_issues"  ]'
add_textbox(s2, RX2 + 0.5, json_y + 0.65, 3.3, 0.2,
            "allowed_tools:", font_size=9, bold=True, color=BODY_COLOR)
add_textbox(s2, RX2 + 0.5, json_y + 0.85, 3.3, 0.25,
            json_text_a, font_size=9, color=GREEN_DARK, font_name="Consolas")
add_textbox(s2, RX2 + 0.5, json_y + 1.1, 3.3, 0.2,
            "→ Sees both Jira and Confluence tools",
            font_size=9, bold=True, color=GREEN_DARK)

# Freemium user entry
add_rounded_rect(s2, RX2 + 4.2, json_y + 0.38, 3.7, 1.1, ORANGE_BG, ORANGE_ACCENT, Pt(1))
add_textbox(s2, RX2 + 4.35, json_y + 0.42, 3.3, 0.22,
            "User B (sub: 557058:d665...)", font_size=10, bold=True, color=ORANGE_DARK)
json_text_b = '[  "jira_list_issues"  ]'
add_textbox(s2, RX2 + 4.35, json_y + 0.65, 3.3, 0.2,
            "allowed_tools:", font_size=9, bold=True, color=BODY_COLOR)
add_textbox(s2, RX2 + 4.35, json_y + 0.85, 3.3, 0.25,
            json_text_b, font_size=9, color=ORANGE_DARK, font_name="Consolas")
add_textbox(s2, RX2 + 4.35, json_y + 1.1, 3.3, 0.2,
            "→ Only Jira tool visible; Confluence blocked",
            font_size=9, bold=True, color=ORANGE_DARK)

# ── Deployment info ──
dep_y = 5.0
dep_card = add_rounded_rect(s2, RX2 + 0.2, dep_y, 7.9, 2.05, WHITE, GRAY_LIGHT, Pt(1))
dep_card.adjustments[0] = 0.04

add_textbox(s2, RX2 + 0.35, dep_y + 0.08, 3.0, 0.25,
            "Deployment & Configuration", font_size=13, bold=True, color=TITLE_COLOR)

dep_items = [
    ("MCP Server", "FastMCP on Azure Container Apps, Streamable HTTP, port 8000"),
    ("Backend API", "FastAPI on ACA, validates JWT via JWKS (Entra ID), exchanges OBO tokens"),
    ("Frontend", "React SPA with MSAL.js, authenticates via Entra ID, acquires scoped tokens"),
    ("Foundry Agent", "gpt-4.1-mini, connected to MCP tool + Fabric Data Agent MCP tool"),
    ("Fabric Integration", "Fabric Data Agent with Row-Level Security on Lakehouse data"),
    ("Agent Instructions", '"Use tools to answer. State which tool. Format URLs as markdown links."'),
]
y = dep_y + 0.38
for label, desc in dep_items:
    tf = add_rich_textbox(s2, RX2 + 0.35, y, 7.4, 0.28)
    p = tf.paragraphs[0]
    rb = p.add_run(); rb.text = label + ": "
    rb.font.size = Pt(10); rb.font.bold = True
    rb.font.color.rgb = BOLD_COLOR; rb.font.name = FONT_NAME
    rd = p.add_run(); rd.text = desc
    rd.font.size = Pt(9); rd.font.color.rgb = BODY_COLOR; rd.font.name = FONT_NAME
    y += 0.27


# ── Save ──
# Try original path; if locked, save with suffix
try:
    prs.save(pptx_path)
    print(f"Saved → {pptx_path}")
except PermissionError:
    alt_path = pptx_path.replace(".pptx", "_updated.pptx")
    prs.save(alt_path)
    print(f"Original locked — saved → {alt_path}")
print(f"Total slides: {len(prs.slides)}")
