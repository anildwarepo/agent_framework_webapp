"""Generate slides: Streaming Use Cases + MCP Inspector Quick Demo."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Colors (matching SVG theme) ──────────────────────────────────────────────
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_BG       = RGBColor(0xF8, 0xFA, 0xFC)
TITLE_COLOR    = RGBColor(0x1E, 0x29, 0x3B)  # dark slate
SUBTITLE_COLOR = RGBColor(0x33, 0x41, 0x55)
BODY_COLOR     = RGBColor(0x47, 0x55, 0x69)
MUTED_COLOR    = RGBColor(0x64, 0x74, 0x8B)
BOLD_COLOR     = RGBColor(0x1E, 0x29, 0x3B)
BLUE_ACCENT    = RGBColor(0x25, 0x63, 0xEB)  # Agent Backend
BLUE_DARK      = RGBColor(0x1E, 0x40, 0xAF)
BLUE_BG        = RGBColor(0xEF, 0xF6, 0xFF)
BLUE_CHIP_BG   = RGBColor(0xDB, 0xEA, 0xFE)
BLUE_CHIP_BD   = RGBColor(0x93, 0xC5, 0xFD)
GREEN_ACCENT   = RGBColor(0x16, 0xA3, 0x4A)  # MCP Server
GREEN_DARK     = RGBColor(0x16, 0x65, 0x34)
GREEN_BG       = RGBColor(0xF0, 0xFD, 0xF4)
GREEN_CHIP_BG  = RGBColor(0xDC, 0xFC, 0xE7)
GREEN_CHIP_BD  = RGBColor(0x86, 0xEF, 0xAC)
PURPLE_ACCENT  = RGBColor(0xA8, 0x55, 0xF7)  # Backend Services
PURPLE_DARK    = RGBColor(0x6B, 0x21, 0xA8)
PURPLE_BG      = RGBColor(0xFA, 0xF5, 0xFF)
GRAY_BD        = RGBColor(0x94, 0xA3, 0xB8)
GRAY_LIGHT     = RGBColor(0xCB, 0xD5, 0xE1)
GRAY_DARK      = RGBColor(0x33, 0x41, 0x55)
YELLOW_BG      = RGBColor(0xFE, 0xF9, 0xC3)
YELLOW_BD      = RGBColor(0xFA, 0xCC, 0x15)
YELLOW_TEXT     = RGBColor(0x85, 0x4D, 0x0E)
LIGHT_BLUE_BG  = RGBColor(0xE0, 0xF2, 0xFE)
LIGHT_BLUE_BD  = RGBColor(0x7D, 0xD3, 0xFC)
CYAN_TEXT       = RGBColor(0x03, 0x69, 0xA1)
ORANGE_ACCENT  = RGBColor(0xEA, 0x58, 0x0C)
ORANGE_DARK    = RGBColor(0x9A, 0x34, 0x12)
ORANGE_BG      = RGBColor(0xFF, 0xF7, 0xED)
ORANGE_CHIP_BG = RGBColor(0xFE, 0xD7, 0xAA)
ORANGE_CHIP_BD = RGBColor(0xFD, 0xBA, 0x74)
DARK_BG        = RGBColor(0x1E, 0x29, 0x3B)
CODE_BG        = RGBColor(0xF1, 0xF5, 0xF9)
CODE_BD        = RGBColor(0xCB, 0xD5, 0xE1)

FONT_NAME = "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)   # widescreen 16:9
prs.slide_height = Inches(7.5)

# ── Helpers (slide-parameterized) ────────────────────────────────────────────

def make_slide_bg(s):
    bg = s.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = SLIDE_BG


def add_textbox(s, left, top, width, height, text, font_size=11, bold=False,
                color=BODY_COLOR, alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def add_rounded_rect(s, left, top, width, height, fill_color, border_color,
                     border_width=Pt(1.5)):
    shape = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    shape.adjustments[0] = 0.08
    return shape


def add_rich_textbox(s, left, top, width, height):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def add_connector(s, x1, y1, x2, y2, color, width=Pt(1.5), dashed=False):
    cx = s.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    cx.line.color.rgb = color
    cx.line.width = width
    if dashed:
        cx.line.dash_style = 4
    return cx


def add_arrow_line(s, x1, y1, x2, y2, color, width=Pt(1.5), dashed=False):
    cx = add_connector(s, x1, y1, x2, y2, color, width, dashed)
    ln = cx.line._ln
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return cx


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Real-Time Streaming Agent Use Cases
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
make_slide_bg(slide)

# Slide title
add_textbox(slide, 0.5, 0.3, 12, 0.55,
            "Why streaming matters for agentic workloads",
            font_size=24, bold=True, color=TITLE_COLOR)

# Use cases card background
card = add_rounded_rect(slide, 0.4, 1.0, 5.6, 5.6, WHITE, GRAY_LIGHT, Pt(1))
card.adjustments[0] = 0.04

# "Use Cases" header
add_textbox(slide, 0.7, 1.15, 4.5, 0.45,
            "Use Cases", font_size=18, bold=True, color=BLUE_ACCENT)

# Bullet items
use_cases = [
    ("Long-running tool calls", " — Database queries, API calls, or file processing that take seconds to minutes; push progress % as work proceeds"),
    ("Multi-step tool chains", " — Agent calls tool A then tool B; surface intermediate results immediately so the user isn't staring at a spinner"),
    ("Human-in-the-loop", " — Push \"I'm about to perform a destructive action — confirm?\" before the tool completes"),
    ("Live status dashboards", " — Show which tool is active, what stage it's in, and estimated time to completion"),
]

y_pos = 1.7
for label, desc in use_cases:
    tf = add_rich_textbox(slide, 0.7, y_pos, 5.0, 1.1)
    p = tf.paragraphs[0]
    p.space_after = Pt(4)

    r_bold = p.add_run()
    r_bold.text = label
    r_bold.font.size = Pt(13)
    r_bold.font.bold = True
    r_bold.font.color.rgb = BOLD_COLOR
    r_bold.font.name = FONT_NAME

    r_desc = p.add_run()
    r_desc.text = desc
    r_desc.font.size = Pt(12)
    r_desc.font.bold = False
    r_desc.font.color.rgb = BODY_COLOR
    r_desc.font.name = FONT_NAME

    y_pos += 1.1


# ── RIGHT SIDE — Architecture Diagram ──

DX = 6.4
DY = 1.0

diag_card = add_rounded_rect(slide, DX - 0.1, DY, 6.7, 5.6, WHITE, GRAY_LIGHT, Pt(1))
diag_card.adjustments[0] = 0.04

add_textbox(slide, DX, DY + 0.15, 6.4, 0.4,
            "Streaming Architecture — Agent Backend → MCP Server → Backend Services",
            font_size=12, bold=True, color=TITLE_COLOR, alignment=PP_ALIGN.CENTER)

bx, by = DX + 0.1, DY + 1.5
bw, bh = 1.35, 1.75
add_rounded_rect(slide, bx, by, bw, bh, WHITE, GRAY_BD, Pt(1.5))
add_textbox(slide, bx, by + 0.1, bw, 0.3, "Browser / Client",
            font_size=10, bold=True, color=SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)
for i, txt in enumerate(["Chat UI", "SSE listener", "Progress indicators"]):
    add_textbox(slide, bx, by + 0.4 + i * 0.25, bw, 0.25, txt,
                font_size=9, color=MUTED_COLOR, alignment=PP_ALIGN.CENTER)
chip = add_rounded_rect(slide, bx + 0.15, by + 1.25, 1.05, 0.3, LIGHT_BLUE_BG, LIGHT_BLUE_BD, Pt(1))
chip.adjustments[0] = 0.15
add_textbox(slide, bx + 0.15, by + 1.27, 1.05, 0.28, "GET /events (SSE)",
            font_size=8, color=CYAN_TEXT, alignment=PP_ALIGN.CENTER)

ax, ay = DX + 1.8, DY + 0.7
aw, ah = 1.6, 3.6
add_rounded_rect(slide, ax, ay, aw, ah, BLUE_BG, BLUE_ACCENT, Pt(2))
add_textbox(slide, ax, ay + 0.08, aw, 0.3, "Agent Backend",
            font_size=11, bold=True, color=BLUE_DARK, alignment=PP_ALIGN.CENTER)

agent_chips = ["LLM / Orchestrator", "MCP Client", "SSE Bus (Redis)", "Session Manager"]
for i, txt in enumerate(agent_chips):
    cy = ay + 0.45 + i * 0.45
    chip = add_rounded_rect(slide, ax + 0.12, cy, 1.36, 0.35, BLUE_CHIP_BG, BLUE_CHIP_BD, Pt(1))
    chip.adjustments[0] = 0.15
    add_textbox(slide, ax + 0.12, cy + 0.02, 1.36, 0.33, txt,
                font_size=9, color=BLUE_DARK, alignment=PP_ALIGN.CENTER)

ycA = ay + 2.45
add_rounded_rect(slide, ax + 0.05, ycA, 1.5, 0.6, YELLOW_BG, YELLOW_BD, Pt(1.5))
add_textbox(slide, ax + 0.05, ycA + 0.05, 1.5, 0.25, "Real-time streaming",
            font_size=8, bold=True, color=YELLOW_TEXT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, ax + 0.05, ycA + 0.28, 1.5, 0.25, "progress → SSE → browser",
            font_size=8, color=YELLOW_TEXT, alignment=PP_ALIGN.CENTER)

mx, my = DX + 3.75, DY + 0.7
mw, mh = 1.45, 3.6
add_rounded_rect(slide, mx, my, mw, mh, GREEN_BG, GREEN_ACCENT, Pt(2))
add_textbox(slide, mx, my + 0.08, mw, 0.3, "MCP Server",
            font_size=11, bold=True, color=GREEN_DARK, alignment=PP_ALIGN.CENTER)

mcp_chips = ["POST /mcp router", "@tool registry", "JSON-RPC 2.0", "Notifications"]
for i, txt in enumerate(mcp_chips):
    cy = my + 0.45 + i * 0.45
    chip = add_rounded_rect(slide, mx + 0.1, cy, 1.25, 0.35, GREEN_CHIP_BG, GREEN_CHIP_BD, Pt(1))
    chip.adjustments[0] = 0.15
    add_textbox(slide, mx + 0.1, cy + 0.02, 1.25, 0.33, txt,
                font_size=9, color=GREEN_DARK, alignment=PP_ALIGN.CENTER)

ycM = my + 2.45
add_rounded_rect(slide, mx + 0.025, ycM, 1.4, 0.6, YELLOW_BG, YELLOW_BD, Pt(1.5))
add_textbox(slide, mx + 0.025, ycM + 0.05, 1.4, 0.25, "ctx.info() / progress",
            font_size=8, bold=True, color=YELLOW_TEXT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, mx + 0.025, ycM + 0.28, 1.4, 0.25, "→ SSE notifications",
            font_size=8, color=YELLOW_TEXT, alignment=PP_ALIGN.CENTER)

sx, sw, sh = DX + 5.55, 1.0, 0.55
services = ["Databases", "APIs / Search", "File Storage", "Auth / Identity"]
for i, svc in enumerate(services):
    sy = DY + 0.8 + i * 0.75
    add_rounded_rect(slide, sx, sy, sw, sh, PURPLE_BG, PURPLE_ACCENT, Pt(1.5))
    add_textbox(slide, sx, sy + 0.08, sw, 0.4, svc,
                font_size=9, bold=True, color=PURPLE_DARK, alignment=PP_ALIGN.CENTER)

add_arrow_line(slide, bx + bw, by + bh / 2, ax, ay + ah / 2 - 0.15, GRAY_BD, Pt(1.5))
add_textbox(slide, bx + bw + 0.02, by + bh / 2 - 0.25, 0.4, 0.2, "HTTP\nSSE",
            font_size=7, color=MUTED_COLOR, alignment=PP_ALIGN.CENTER)

arr_y1 = ay + 1.15
add_arrow_line(slide, ax + aw, arr_y1, mx, arr_y1, BLUE_ACCENT, Pt(2))
add_textbox(slide, ax + aw + 0.02, arr_y1 - 0.22, 0.6, 0.2, "tools/call",
            font_size=8, bold=True, color=BLUE_ACCENT, alignment=PP_ALIGN.CENTER)

arr_y2 = ay + 1.65
add_arrow_line(slide, mx, arr_y2, ax + aw, arr_y2, GREEN_ACCENT, Pt(2), dashed=True)
add_textbox(slide, ax + aw + 0.02, arr_y2 + 0.03, 0.8, 0.2, "SSE notifications",
            font_size=8, bold=True, color=GREEN_ACCENT, alignment=PP_ALIGN.CENTER)

for i in range(4):
    svc_y = DY + 0.8 + i * 0.75 + sh / 2
    mcp_y = my + 0.45 + i * 0.45 + 0.175
    add_arrow_line(slide, mx + mw, mcp_y, sx, svc_y, PURPLE_ACCENT, Pt(1.5))

leg_y = DY + 4.7
add_rounded_rect(slide, DX + 0.1, leg_y, 6.4, 0.55, WHITE, GRAY_LIGHT, Pt(1))

items = [
    (DX + 0.3, "Request (tools/call)", BLUE_ACCENT, False),
    (DX + 2.0, "SSE streaming notifications", GREEN_ACCENT, True),
    (DX + 4.2, "Backend service calls", PURPLE_ACCENT, False),
]
for lx, label, color, dashed in items:
    add_connector(slide, lx, leg_y + 0.275, lx + 0.35, leg_y + 0.275, color, Pt(2), dashed)
    add_textbox(slide, lx + 0.4, leg_y + 0.12, 1.6, 0.3, label,
                font_size=9, color=SUBTITLE_COLOR)

add_rounded_rect(slide, DX + 5.65, leg_y + 0.15, 0.18, 0.18, YELLOW_BG, YELLOW_BD, Pt(1))
add_textbox(slide, DX + 5.88, leg_y + 0.12, 0.9, 0.3, "Streaming callout",
            font_size=9, color=SUBTITLE_COLOR)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Quick Demo: MCP Inspector
# ══════════════════════════════════════════════════════════════════════════════

s2 = prs.slides.add_slide(prs.slide_layouts[6])
make_slide_bg(s2)

# ── Title ──
add_textbox(s2, 0.5, 0.3, 12, 0.55,
            "Quick Demo: MCP Inspector",
            font_size=24, bold=True, color=TITLE_COLOR)
add_textbox(s2, 0.5, 0.85, 12, 0.35,
            "Interactive debugging & testing tool for MCP servers — no agent required",
            font_size=14, color=MUTED_COLOR)

# ══════════════════════════════════════════════════════════════════════════════
# LEFT SIDE — What / How / Demo steps
# ══════════════════════════════════════════════════════════════════════════════

LX = 0.4
card2 = add_rounded_rect(s2, LX, 1.4, 5.8, 5.4, WHITE, GRAY_LIGHT, Pt(1))
card2.adjustments[0] = 0.04

add_textbox(s2, LX + 0.3, 1.55, 5.2, 0.35,
            "What is MCP Inspector?", font_size=16, bold=True, color=BLUE_ACCENT)

what_items = [
    ("Browser-based UI", " that connects to any MCP server (stdio or Streamable HTTP)"),
    ("Discover tools", " via initialize → tools/list — see names, schemas, descriptions"),
    ("Call tools interactively", " — fill in args, fire tools/call, inspect JSON-RPC responses"),
    ("Watch SSE notifications", " — progress, messages, and logging events stream live"),
]
y = 1.95
for label, desc in what_items:
    tf = add_rich_textbox(s2, LX + 0.3, y, 5.1, 0.4)
    p = tf.paragraphs[0]
    rb = p.add_run(); rb.text = "• " + label
    rb.font.size = Pt(11); rb.font.bold = True
    rb.font.color.rgb = BOLD_COLOR; rb.font.name = FONT_NAME
    rd = p.add_run(); rd.text = desc
    rd.font.size = Pt(11); rd.font.color.rgb = BODY_COLOR; rd.font.name = FONT_NAME
    y += 0.4

# Divider
add_connector(s2, LX + 0.3, y + 0.1, LX + 5.5, y + 0.1, GRAY_LIGHT, Pt(1))

# "Demo Steps" section
add_textbox(s2, LX + 0.3, y + 0.25, 5.2, 0.35,
            "Demo Steps", font_size=16, bold=True, color=GREEN_ACCENT)

demo_steps = [
    ("1.", "Launch", "npx @modelcontextprotocol/inspector"),
    ("2.", "Connect", "Enter MCP server URL → http://localhost:3000/mcp"),
    ("3.", "Discover", "Click \"List Tools\" — see all registered @tool functions"),
    ("4.", "Execute", "Pick a tool → fill args → \"Call Tool\" → see JSON result"),
    ("5.", "Stream", "Watch SSE panel — progress %, ctx.info() messages live"),
]
y2 = y + 0.65
for num, label, desc in demo_steps:
    tf = add_rich_textbox(s2, LX + 0.3, y2, 5.1, 0.42)
    p = tf.paragraphs[0]
    rn = p.add_run(); rn.text = num + " "
    rn.font.size = Pt(12); rn.font.bold = True
    rn.font.color.rgb = GREEN_DARK; rn.font.name = FONT_NAME
    rl = p.add_run(); rl.text = label + "  "
    rl.font.size = Pt(12); rl.font.bold = True
    rl.font.color.rgb = BOLD_COLOR; rl.font.name = FONT_NAME
    rd = p.add_run(); rd.text = desc
    rd.font.size = Pt(11); rd.font.color.rgb = BODY_COLOR; rd.font.name = FONT_NAME
    y2 += 0.42

# Command chip
cmd_y = y2 + 0.15
cmd = add_rounded_rect(s2, LX + 0.3, cmd_y, 5.2, 0.4, CODE_BG, CODE_BD, Pt(1))
cmd.adjustments[0] = 0.15
add_textbox(s2, LX + 0.5, cmd_y + 0.04, 4.8, 0.32,
            "$ npx @modelcontextprotocol/inspector --url http://localhost:3000/mcp",
            font_size=10, color=TITLE_COLOR, font_name="Consolas")

# ══════════════════════════════════════════════════════════════════════════════
# RIGHT SIDE — MCP Inspector Architecture (native shapes)
# ══════════════════════════════════════════════════════════════════════════════

RX = 6.6
diag2 = add_rounded_rect(s2, RX - 0.1, 1.4, 6.5, 5.4, WHITE, GRAY_LIGHT, Pt(1))
diag2.adjustments[0] = 0.04

add_textbox(s2, RX, 1.55, 6.2, 0.35,
            "MCP Inspector — How It Works",
            font_size=14, bold=True, color=TITLE_COLOR, alignment=PP_ALIGN.CENTER)

# ── Developer box ──
dx2, dy2 = RX + 0.3, 2.2
dw2, dh2 = 1.4, 1.2
add_rounded_rect(s2, dx2, dy2, dw2, dh2, WHITE, GRAY_BD, Pt(1.5))
add_textbox(s2, dx2, dy2 + 0.12, dw2, 0.3, "Developer",
            font_size=11, bold=True, color=SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)
add_textbox(s2, dx2, dy2 + 0.42, dw2, 0.25, "Browser UI",
            font_size=9, color=MUTED_COLOR, alignment=PP_ALIGN.CENTER)
add_textbox(s2, dx2, dy2 + 0.65, dw2, 0.25, "localhost:6274",
            font_size=9, color=MUTED_COLOR, alignment=PP_ALIGN.CENTER)
add_textbox(s2, dx2, dy2 + 0.88, dw2, 0.25, "Fill args & run",
            font_size=9, color=MUTED_COLOR, alignment=PP_ALIGN.CENTER)

# ── MCP Inspector box (orange) ──
ix, iy = RX + 2.2, 2.2
iw, ih = 1.9, 2.9
add_rounded_rect(s2, ix, iy, iw, ih, ORANGE_BG, ORANGE_ACCENT, Pt(2))
add_textbox(s2, ix, iy + 0.08, iw, 0.3, "MCP Inspector",
            font_size=12, bold=True, color=ORANGE_DARK, alignment=PP_ALIGN.CENTER)

insp_chips = ["initialize", "tools/list", "tools/call", "SSE viewer", "JSON-RPC log"]
for i, txt in enumerate(insp_chips):
    cy = iy + 0.5 + i * 0.45
    chip = add_rounded_rect(s2, ix + 0.15, cy, 1.6, 0.35, ORANGE_CHIP_BG, ORANGE_CHIP_BD, Pt(1))
    chip.adjustments[0] = 0.15
    add_textbox(s2, ix + 0.15, cy + 0.02, 1.6, 0.33, txt,
                font_size=9, color=ORANGE_DARK, alignment=PP_ALIGN.CENTER)

# ── MCP Server box (green) ──
sx2, sy2_top = RX + 4.6, 2.2
sw2, sh2 = 1.5, 2.9
add_rounded_rect(s2, sx2, sy2_top, sw2, sh2, GREEN_BG, GREEN_ACCENT, Pt(2))
add_textbox(s2, sx2, sy2_top + 0.08, sw2, 0.3, "MCP Server",
            font_size=12, bold=True, color=GREEN_DARK, alignment=PP_ALIGN.CENTER)

srv_chips = ["POST /mcp", "@tool functions", "Notifications", "Backend calls"]
for i, txt in enumerate(srv_chips):
    cy = sy2_top + 0.5 + i * 0.55
    chip = add_rounded_rect(s2, sx2 + 0.1, cy, 1.3, 0.35, GREEN_CHIP_BG, GREEN_CHIP_BD, Pt(1))
    chip.adjustments[0] = 0.15
    add_textbox(s2, sx2 + 0.1, cy + 0.02, 1.3, 0.33, txt,
                font_size=9, color=GREEN_DARK, alignment=PP_ALIGN.CENTER)

# ── Arrows ──
add_arrow_line(s2, dx2 + dw2, dy2 + dh2 / 2, ix, iy + 0.6, GRAY_BD, Pt(1.5))
add_textbox(s2, dx2 + dw2 + 0.02, dy2 + dh2 / 2 - 0.28, 0.7, 0.2, "click / fill args",
            font_size=7, color=MUTED_COLOR, alignment=PP_ALIGN.CENTER)

add_arrow_line(s2, ix + iw, iy + 0.9, sx2, sy2_top + 0.9, ORANGE_ACCENT, Pt(2))
add_textbox(s2, ix + iw + 0.02, iy + 0.65, 0.65, 0.2, "JSON-RPC",
            font_size=8, bold=True, color=ORANGE_ACCENT, alignment=PP_ALIGN.CENTER)

add_arrow_line(s2, sx2, sy2_top + 1.6, ix + iw, iy + 1.6, GREEN_ACCENT, Pt(2), dashed=True)
add_textbox(s2, ix + iw + 0.02, iy + 1.68, 0.65, 0.2, "SSE events",
            font_size=8, bold=True, color=GREEN_ACCENT, alignment=PP_ALIGN.CENTER)

# ── Callout boxes ──
cb_y = 5.5

add_rounded_rect(s2, RX + 0.2, cb_y, 2.8, 1.1, YELLOW_BG, YELLOW_BD, Pt(1.5))
add_textbox(s2, RX + 0.3, cb_y + 0.06, 2.6, 0.25, "What you see in the Inspector",
            font_size=9, bold=True, color=YELLOW_TEXT, alignment=PP_ALIGN.CENTER)
see_items = ["Tool list + input schemas", "JSON-RPC request / response", "Live SSE notification stream"]
sy_i = cb_y + 0.3
for item in see_items:
    add_textbox(s2, RX + 0.35, sy_i, 2.5, 0.22, "• " + item,
                font_size=8, color=YELLOW_TEXT)
    sy_i += 0.22

add_rounded_rect(s2, RX + 3.3, cb_y, 2.9, 1.1, BLUE_CHIP_BG, BLUE_CHIP_BD, Pt(1.5))
add_textbox(s2, RX + 3.4, cb_y + 0.06, 2.7, 0.25, "Why it matters for demos",
            font_size=9, bold=True, color=BLUE_DARK, alignment=PP_ALIGN.CENTER)
why_items = ["No agent code needed to test tools", "Validate streaming before integration", "Debug JSON-RPC on the wire"]
wy_i = cb_y + 0.3
for item in why_items:
    add_textbox(s2, RX + 3.45, wy_i, 2.6, 0.22, "• " + item,
                font_size=8, color=BLUE_DARK)
    wy_i += 0.22

# ── Save ──
out_path = os.path.join(os.path.dirname(__file__), "mcp_slides.pptx")
prs.save(out_path)
print(f"Saved → {out_path}")
